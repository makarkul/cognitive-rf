"""Build design-discussion summary PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Palette: Ocean Gradient
DEEP = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
MID = RGBColor(0x21, 0x29, 0x5C)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x5F, 0x6D)
ACCENT = RGBColor(0xF5, 0xA6, 0x23)

HDR_FONT = "Calibri"
BODY_FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def add_bg(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, size=16, bold=False,
             color=DARK, font=BODY_FONT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_accent_bar(slide, left, top, width=Inches(0.12), height=Inches(0.6), color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def slide_title(slide, title, subtitle=None, dark=False):
    color = LIGHT if dark else DEEP
    add_text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8),
             title, size=30, bold=True, color=color, font=HDR_FONT)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.4),
                 subtitle, size=14, color=MUTED if not dark else RGBColor(0xCC, 0xD6, 0xE4),
                 font=BODY_FONT, italic=False)


def add_table(slide, left, top, width, height, data, header=True,
              col_widths=None, font_size=12, header_color=DEEP, header_text_color=LIGHT):
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(width * w / total)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.text = ""
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.name = BODY_FONT
            run.font.size = Pt(font_size)
            if header and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                run.font.bold = True
                run.font.color.rgb = header_text_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r % 2 == 1 else RGBColor(0xE8, 0xEE, 0xF5)
                run.font.color.rgb = DARK
    return tbl


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(blank)
add_bg(s, MID)
# Side accent band

add_text(s, Inches(1.0), Inches(2.2), Inches(11.5), Inches(1.6),
         "Sinusoid Recovery with a Tiny Transformer",
         size=44, bold=True, color=LIGHT, font=HDR_FONT)
add_text(s, Inches(1.0), Inches(3.7), Inches(11.5), Inches(0.8),
         "A DSP-to-Transformer Walkthrough",
         size=24, color=RGBColor(0xCA, 0xDC, 0xFC), font=HDR_FONT, italic=False)
add_text(s, Inches(1.0), Inches(5.6), Inches(11), Inches(0.5),
         "Design notes: AR modeling, attention, and channel-coding intuitions",
         size=14, color=RGBColor(0x97, 0xAC, 0xC7), font=BODY_FONT)

# ---------- Slide 2: Motivation ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Motivation", "Classical DSP and transformers — the same job?")

left_col = Inches(0.85); col_w = Inches(5.9)
add_text(s, left_col, Inches(1.7), col_w, Inches(0.5),
         "Classical DSP view", size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, left_col, Inches(2.2), col_w, Inches(4),
         "• Fit AR / MA / ARMA models to a time series\n"
         "• PACF determines the AR order p\n"
         "• Coefficients fit by Yule-Walker or least squares\n"
         "• Prediction = linear combination of past samples",
         size=15)

add_text(s, Inches(7.0), Inches(1.7), col_w, Inches(0.5),
         "Transformer view", size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(7.0), Inches(2.2), col_w, Inches(4),
         "• Train to predict next sample from a context\n"
         "• Self-attention weights pick which lags matter\n"
         "• Coefficients learned via gradient descent\n"
         "• Thesis: attention should rediscover AR structure",
         size=15)

add_text(s, Inches(0.85), Inches(6.1), Inches(11.5), Inches(0.7),
         "If the thesis holds, a tiny transformer is just an AR model in disguise.",
         size=16, italic=False, color=DEEP, font=HDR_FONT)

# ---------- Slide 3: Key DSP fact ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Key DSP Fact", "A sinusoid is exactly an AR(2) process")

eq_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.3),
                             Inches(10.3), Inches(1.6))
eq_box.fill.solid(); eq_box.fill.fore_color.rgb = DEEP; eq_box.line.fill.background()
add_text(s, Inches(1.5), Inches(2.55), Inches(10.3), Inches(1.1),
         "x[n] = 2 cos(w0) · x[n-1]  −  x[n-2]",
         size=32, bold=True, color=LIGHT, font=BODY_FONT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.85), Inches(4.3), Inches(11.5), Inches(0.5),
         "What this means",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(4.8), Inches(11.5), Inches(2.5),
         "• Only two past samples carry information — lag-1 and lag-2\n"
         "• The single coefficient  2 cos(w0)  encodes frequency\n"
         "• The model needs only 2 effective taps, regardless of freq\n"
         "• This sets a hard lower bound on how tiny the network can be",
         size=16)

# ---------- Slide 4: Problem setup ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Problem Setup", "Noisy sinusoid → clean next-sample prediction")

data = [
    ["Parameter", "Value / Range", "Role"],
    ["Frequency", "1 – 20 Hz (uniform)", "Randomized per sequence"],
    ["Amplitude", "0.5 – 2.0 (uniform)", "Randomized per sequence"],
    ["Phase", "0 – 2π (uniform)", "Randomized per sequence"],
    ["Sampling rate fs", "100 Hz", "Fixed, well above Nyquist"],
    ["Amplitude noise", "σ = 0.2 (Gaussian)", "Corrupts amplitude"],
    ["Phase noise", "σ = 0.1 rad (Gaussian)", "Corrupts phase"],
    ["Context length", "128 samples", "≈ 1 period at f_min = 1 Hz"],
]
add_table(s, Inches(0.85), Inches(1.9), Inches(11.6), Inches(4.4), data,
          col_widths=[3, 4, 5], font_size=13)

add_text(s, Inches(0.85), Inches(6.5), Inches(11.5), Inches(0.7),
         "Task: given 128 noisy samples, predict the next clean sample — at every position.",
         size=15, italic=False, color=DEEP, font=HDR_FONT)

# ---------- Slide 5: Architecture ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Architecture", "Scalar in → embed → 2 transformer blocks → scalar out")

# Pipeline boxes
stages = [
    ("Scalar\nsample", RGBColor(0xB8, 0xC7, 0xD8)),
    ("Linear(1,32)\ninput proj", TEAL),
    ("+ pos emb\n(128×32)", TEAL),
    ("Trf block ×2\n2 heads, 32 dim", DEEP),
    ("LayerNorm", TEAL),
    ("Linear(32,1)\nout head", TEAL),
    ("Next\nsample", ACCENT),
]
x = Inches(0.45); y = Inches(2.5); w = Inches(1.65); h = Inches(1.3); gap = Inches(0.18)
for i, (label, color) in enumerate(stages):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
    text_color = LIGHT if color in (TEAL, DEEP, ACCENT) else DARK
    tf = box.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for j, line in enumerate(label.split("\n")):
        if j == 0:
            r = p.add_run()
        else:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r = p2.add_run()
        r.text = line
        r.font.name = BODY_FONT; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = text_color
    if i < len(stages) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    x + w, y + Inches(0.5), gap, Inches(0.3))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED; arrow.line.fill.background()
    x += w + gap

add_text(s, Inches(0.85), Inches(4.5), Inches(11.5), Inches(0.5),
         "~29,000 parameters total",
         size=20, bold=True, color=DEEP, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(5.1), Inches(11.5), Inches(2),
         "• Differences from GPT: nn.Embedding → Linear(1, 32); softmax head → Linear(32, 1)\n"
         "• Reused from ch04: MultiHeadAttention, TransformerBlock, LayerNorm, GELU, FeedForward\n"
         "• Loss: MSE instead of cross-entropy — the target is a real value, not a class",
         size=15)

# ---------- Slide 6: TL;DR table ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "The TL;DR Design Table", "Three DSP-anchored hyperparameters drive everything")

data = [
    ["Hyperparameter", "What it is", "DSP anchor", "Sizing rule", "This project"],
    ["n_heads", "Parallel attention patterns", "Number of AR taps / lags", "≈ effective AR order", "2"],
    ["context_length", "Lookback window", "N_FFT — covers slowest period", "≈ fs / f_min", "128"],
    ["head_dim", "Q·K discrimination width", "Codeword length for positions", "≈ 2 · log₂(context_length)", "16"],
]
add_table(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(3.2), data,
          col_widths=[2, 3, 3.5, 3.5, 1.5], font_size=12)

callout = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(5.4),
                              Inches(11.6), Inches(1.5))
callout.fill.solid(); callout.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF5)
callout.line.color.rgb = ACCENT; callout.line.width = Pt(1.5)
add_text(s, Inches(1.0), Inches(5.55), Inches(11.3), Inches(0.5),
         "emb_dim is not a free knob",
         size=16, bold=True, color=DEEP, font=HDR_FONT)
add_text(s, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.9),
         "emb_dim = n_heads × head_dim = 2 × 16 = 32.  It falls out automatically from the three choices above.",
         size=14, color=DARK)

# ---------- Slide 7: Tokenization ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Tokenization Is Different", "No tokenizer, no vocab_size")

data = [
    ["Aspect", "GPT (text)", "This model (signals)"],
    ["Input", "Discrete token IDs", "Continuous scalar samples"],
    ["Embedding", "nn.Embedding(vocab_size, 32)", "nn.Linear(1, 32)"],
    ["Vocabulary", "~50,000 tokens", "Does not exist"],
    ["Output head", "Linear(32, vocab_size) + softmax", "Linear(32, 1)"],
    ["Loss", "Cross-entropy", "MSE"],
]
add_table(s, Inches(0.85), Inches(1.9), Inches(11.6), Inches(3.6), data,
          col_widths=[2, 4, 4.5], font_size=13)

add_text(s, Inches(0.85), Inches(5.8), Inches(11.5), Inches(1.5),
         "The embedding layer's only job: lift a scalar into the 32-D working space\n"
         "where positional embeddings and attention machinery can operate on it.",
         size=15, italic=False, color=MUTED)

# ---------- Slide 8: Is 32-D wasted? ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Are the 32 Dimensions Wasted?", "Why a scalar gets a 32-D embedding")

add_text(s, Inches(0.85), Inches(1.9), Inches(11.6), Inches(0.6),
         "Linear(1, 32) produces  e = v · w  for scalar value v and a fixed 32-D vector w.",
         size=15, color=DARK, font=BODY_FONT)

add_text(s, Inches(0.85), Inches(2.7), Inches(11.6), Inches(0.5),
         "The signal content is rank-1",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(3.2), Inches(11.6), Inches(1.5),
         "• All 32 dimensions are perfectly correlated in v\n"
         "• Information content about the sample itself: still 1 scalar\n"
         "• No magic expansion — the embedding just places v on a line in 32-D space",
         size=14)

add_text(s, Inches(0.85), Inches(4.9), Inches(11.6), Inches(0.5),
         "So why 32? For the rest of the machinery",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(1.7),
         "• Positional embeddings live in the same 32-D space — position needs real capacity\n"
         "• Attention's Q/K dot products need room for near-orthogonal patterns\n"
         "• FFN nonlinearity mixes value and position once they share dimensions",
         size=14)

# ---------- Slide 9: N_FFT analogy ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "context_length ≈ N_FFT", "Cover the slowest period you care about")

add_text(s, Inches(0.85), Inches(1.85), Inches(11.6), Inches(0.5),
         "Invariant:  N_FFT / fs  ≈  1 / f_min",
         size=18, bold=True, color=DEEP, font=BODY_FONT)

data = [
    ["System", "Bandwidth", "fs (sampling)", "N_FFT"],
    ["LTE", "1.4 MHz", "1.92 MHz", "128"],
    ["LTE", "5 MHz", "7.68 MHz", "512"],
    ["LTE", "10 MHz", "15.36 MHz", "1024"],
    ["LTE", "20 MHz", "30.72 MHz", "2048"],
    ["This project", "1 – 20 Hz", "100 Hz", "128"],
]
add_table(s, Inches(0.85), Inches(2.7), Inches(11.6), Inches(3.5), data,
          col_widths=[3, 3, 3, 2], font_size=13)

add_text(s, Inches(0.85), Inches(6.4), Inches(11.6), Inches(0.7),
         "Same idea: pick the window that captures at least one full cycle of the lowest frequency.",
         size=14, italic=False, color=MUTED)

# ---------- Slide 10: n_heads = AR order ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "n_heads = Lags That Matter", "One head per pattern to attend to")

add_text(s, Inches(0.85), Inches(1.9), Inches(11.6), Inches(0.5),
         "Heuristic",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(2.4), Inches(11.6), Inches(1.3),
         "• For AR(p), you want roughly p heads\n"
         "• Each head can specialize on one lag pattern",
         size=15)

add_text(s, Inches(0.85), Inches(4.0), Inches(11.6), Inches(0.5),
         "This project",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(4.5), Inches(11.6), Inches(2.2),
         "• Sinusoid = AR(2) ⇒ n_heads = 2\n"
         "• Head 1 learns the lag-1 pattern (weighted by 2 cos ω₀)\n"
         "• Head 2 learns the lag-2 pattern (constant coefficient −1)\n"
         "• Causal mask ensures strictly past-only attention",
         size=15)

# ---------- Slide 11: head_dim as channel coding ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "head_dim = Channel Coding", "Q·K match = matched-filter detection")

add_text(s, Inches(0.85), Inches(1.85), Inches(11.6), Inches(1.6),
         "head_dim is the dimension of the space Q and K live in.\n"
         "It governs how many distinct positional patterns you can discriminate under noise.",
         size=15)

add_text(s, Inches(0.85), Inches(3.3), Inches(11.6), Inches(0.5),
         "Classical bounds all say the same thing",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(3.8), Inches(11.6), Inches(2.3),
         "• Shannon capacity:     d ≥ 2R / log₂(1 + SNR)\n"
         "• Spherical code packing:  nearly orthogonal vectors need O(log N) dimensions\n"
         "• BCH / Reed-Muller codes:  redundancy ≈ 2 · information bits\n\n"
         "Convergent answer:  head_dim ≈ 2 · log₂(context_length)",
         size=14, font=BODY_FONT)

add_text(s, Inches(0.85), Inches(6.3), Inches(11.6), Inches(0.6),
         "This project: 2 · log₂(128) = 14 ≈ 16",
         size=16, bold=True, color=DEEP, font=HDR_FONT)

# ---------- Slide 12: Channel coding analogy (deeper) ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Why the Analogy Holds", "Scaling, spreading, processing gain")

add_text(s, Inches(0.85), Inches(1.9), Inches(11.6), Inches(0.5),
         "Attention formula",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(2.4), Inches(11.6), Inches(0.7),
         "softmax( Q · Kᵀ / √head_dim ) · V",
         size=16, color=DARK, font=BODY_FONT)

add_text(s, Inches(0.85), Inches(3.3), Inches(11.6), Inches(0.5),
         "Interpretation",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(3.8), Inches(11.6), Inches(3.3),
         "• Each key vector = a template (like a PN spreading code)\n"
         "• Q · K dot product = matched-filter correlation\n"
         "• 1 / √head_dim = processing-gain normalization (variance control)\n"
         "• Longer head_dim  ⇒  less interference between distinct patterns\n"
         "• Softmax = winner-take-all detection after correlation",
         size=14)

# ---------- Slide 13: Shannon sizing ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Shannon Sizing Check", "Does head_dim = 16 meet the information bound?")

data = [
    ["Quantity", "Value", "Note"],
    ["Positions to discriminate", "128", "= context_length"],
    ["Information needed", "7 bits", "= log₂(128)"],
    ["Per-dim SNR (assumed)", "~1", "Training noise level"],
    ["Bits per dim  (½ log₂(1+SNR))", "0.5", "Shannon per-dim capacity"],
    ["Minimum head_dim", "14", "= 2 · 7"],
    ["Chosen head_dim", "16", "Shannon floor × ~1.1"],
    ["Large LLMs use", "~40 – 80", "3 – 5× safety factor"],
]
add_table(s, Inches(1.2), Inches(1.85), Inches(10.9), Inches(4.8), data,
          col_widths=[4, 2, 5], font_size=13)

add_text(s, Inches(0.85), Inches(6.8), Inches(11.6), Inches(0.5),
         "Our 16 is snug — deliberate, not generous. Text models need more headroom.",
         size=14, italic=False, color=MUTED)

# ---------- Slide 14: Training ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Training", "30 epochs on CPU")

add_text(s, Inches(0.85), Inches(1.9), Inches(5.8), Inches(0.5),
         "Recipe", size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(2.4), Inches(5.8), Inches(3.5),
         "• Loss: MSE\n"
         "• Optimizer: AdamW (lr = 1e-3)\n"
         "• Batch size: 32\n"
         "• 10 000 fresh sinusoids / epoch\n"
         "• On-the-fly data generation",
         size=15)

# Results cards
def stat_card(left, top, label, value, sub=None):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(1.5))
    card.fill.solid(); card.fill.fore_color.rgb = DEEP; card.line.fill.background()
    add_text(s, left, top + Inches(0.05), Inches(2.8), Inches(0.4),
             label, size=11, color=RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.CENTER)
    add_text(s, left, top + Inches(0.4), Inches(2.8), Inches(0.8),
             value, size=26, bold=True, color=LIGHT, font=HDR_FONT, align=PP_ALIGN.CENTER)
    if sub:
        add_text(s, left, top + Inches(1.05), Inches(2.8), Inches(0.4),
                 sub, size=10, color=RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.CENTER)

stat_card(Inches(7.0), Inches(1.9), "Initial MSE", "1.40")
stat_card(Inches(10.0), Inches(1.9), "Final MSE", "0.03")
stat_card(Inches(7.0), Inches(3.6), "Output SNR", "14.8 dB", "matches naive baseline")
stat_card(Inches(10.0), Inches(3.6), "Parameters", "29 K")

add_text(s, Inches(0.85), Inches(6.2), Inches(11.6), Inches(0.9),
         "Model learns to predict the next clean sample from noisy history — denoising + AR prediction in one.",
         size=14, italic=False, color=MUTED)

# ---------- Slide 15: Analysis ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Analysis Diagnostics", "Does the model really learn AR(2)?")

items = [
    ("Attention heatmaps", "Peaks at lag 1 and lag 2 — matches AR(2) tap structure"),
    ("Attention-vs-lag profile", "Bar heights track theoretical  2 cos(ω₀)  weights"),
    ("FFT comparison", "Output spectrum concentrates at true frequency, noise suppressed"),
    ("Autoregressive generation", "Seeding with noisy samples yields clean sinusoid continuation"),
]
y = Inches(1.95)
for title_t, desc in items:
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.1),
                                 Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT; circle.line.fill.background()
    add_text(s, Inches(1.5), y, Inches(11), Inches(0.5),
             title_t, size=17, bold=True, color=DEEP, font=HDR_FONT)
    add_text(s, Inches(1.5), y + Inches(0.5), Inches(11), Inches(0.6),
             desc, size=14, color=DARK)
    y += Inches(1.2)

# ---------- Slide 16: Low-SNR sweep ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Performance vs Input SNR", "Original 29 K model, easy-noise training")

data = [
    ["amp_std", "phase_std", "Input SNR (dB)", "Output SNR (dB)", "Gain (dB)"],
    ["0.10", "0.05", "3.05", "17.31", "+14.26"],
    ["0.20", "0.10", "2.85", "14.65", "+11.80  (training)"],
    ["0.40", "0.20", "2.14", "9.61", "+7.47"],
    ["0.60", "0.30", "1.17", "5.90", "+4.72"],
    ["0.80", "0.40", "0.12", "3.26", "+3.14"],
    ["1.00", "0.50", "-0.93", "1.37", "+2.30"],
]
add_table(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(3.8), data,
          col_widths=[2, 2, 3, 3, 3], font_size=13)
add_text(s, Inches(0.85), Inches(6.0), Inches(11.6), Inches(1.0),
         "Gain collapses out-of-distribution. Two effects: harder physical task + train/test mismatch.",
         size=14, color=MUTED)

# ---------- Slide 17: Retrain comparison at hard noise ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Tweaking for Low SNR", "Retrain at hard noise (input SNR ≈ 0 dB)")

data = [
    ["Scenario", "Arch (ctx / emb / heads / layers)", "Params", "Gain (dB)"],
    ["Original (no retrain)", "128 / 32 / 2 / 2", "29 K", "+3.14"],
    ["Same arch, retrained", "128 / 32 / 2 / 2", "29 K", "+3.69"],
    ["Scaled arch, retrained", "256 / 72 / 4 / 3", "208 K", "+6.00"],
]
add_table(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(2.4), data,
          col_widths=[3.5, 4.5, 1.5, 2.5], font_size=13)

add_text(s, Inches(0.85), Inches(4.5), Inches(11.6), Inches(0.5),
         "DSP-grounded tuning rules at low SNR",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(5.0), Inches(11.6), Inches(2.3),
         "• context_length ↑   — extra √k noise averaging\n"
         "• n_heads ↑              — add smoothing/MA heads beyond AR taps\n"
         "• n_layers ↑              — depth for iterative denoising refinement\n"
         "• head_dim follows 2·log₂(context_length); emb_dim falls out\n"
         "• Match training noise to deployment noise; use mixed-σ curriculum",
         size=14)

# ---------- Slide 18: FFT embedding head-to-head ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "FFT Embedding — Head-to-Head", "Replace Linear(1, 32) with sliding-window FFT")

add_text(s, Inches(0.85), Inches(1.9), Inches(11.6), Inches(0.5),
         "Both models ~30 K params; only the input embedding differs",
         size=14, color=MUTED)

# Easy noise table
add_text(s, Inches(0.85), Inches(2.5), Inches(11.6), Inches(0.4),
         "Easy noise  (amp_std=0.2, phase_std=0.1) — 8 epochs",
         size=15, bold=True, color=TEAL, font=HDR_FONT)
data_easy = [
    ["Variant", "Params", "Out SNR", "Gain"],
    ["Scalar Linear(1, 32)", "29,473", "9.92 dB", "+7.31 dB"],
    ["FFT sliding window", "30,529", "16.39 dB", "+13.78 dB"],
]
add_table(s, Inches(0.85), Inches(2.95), Inches(11.6), Inches(1.3), data_easy,
          col_widths=[4, 2, 3, 3], font_size=13)

# Hard noise table
add_text(s, Inches(0.85), Inches(4.5), Inches(11.6), Inches(0.4),
         "Hard noise  (amp_std=0.8, phase_std=0.4) — 10 epochs",
         size=15, bold=True, color=TEAL, font=HDR_FONT)
data_hard = [
    ["Variant", "Params", "Out SNR", "Gain"],
    ["Scalar Linear(1, 32)", "29,473", "3.81 dB", "+3.69 dB"],
    ["FFT sliding window", "30,529", "9.12 dB", "+9.00 dB"],
    ["Scalar SCALED (208 K)", "207,577", "5.85 dB", "+6.00 dB"],
]
add_table(s, Inches(0.85), Inches(4.95), Inches(11.6), Inches(1.7), data_hard,
          col_widths=[4, 2, 3, 3], font_size=13)

add_text(s, Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.5),
         "FFT at 30 K beats the scaled scalar model at 208 K by +3 dB.",
         size=14, bold=True, color=DEEP, font=HDR_FONT)

# ---------- Slide 19: Why FFT helps ----------
s = prs.slides.add_slide(blank); add_bg(s)
slide_title(s, "Why the FFT Front-End Wins", "A near-diagonal representation of the task")

add_text(s, Inches(0.85), Inches(1.9), Inches(11.6), Inches(0.5),
         "Token at position t = Linear( rfft( x[t-W+1 : t+1] ) ),  W = 32",
         size=15, color=DARK, font=BODY_FONT)

add_text(s, Inches(0.85), Inches(2.7), Inches(11.6), Inches(0.5),
         "What changes for the model",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(3.2), Inches(11.6), Inches(2.3),
         "• A pure sinusoid is a delta in the frequency domain\n"
         "• Each token already encodes the local spectrum of the last 32 samples\n"
         "• Noisy input becomes a peak-plus-noise-floor in each token\n"
         "• The transformer's job collapses to peak-tracking and bin interpolation\n"
         "• Equivalent to many parallel matched filters before attention even runs",
         size=14)

add_text(s, Inches(0.85), Inches(5.6), Inches(11.6), Inches(0.5),
         "Trade-offs",
         size=18, bold=True, color=TEAL, font=HDR_FONT)
add_text(s, Inches(0.85), Inches(6.1), Inches(11.6), Inches(1.3),
         "• Inference overhead of sliding FFT is negligible (~160 MACs/sample)\n"
         "• Strong inductive bias — best for stationary/periodic signals\n"
         "• Window size W is a classical STFT trade-off (frequency vs time resolution)",
         size=14)

# ---------- Slide 20: Key takeaway ----------
s = prs.slides.add_slide(blank); add_bg(s, MID)

add_text(s, Inches(1.0), Inches(0.8), Inches(11.5), Inches(1.0),
         "Key Takeaway", size=38, bold=True, color=LIGHT, font=HDR_FONT)
add_text(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(0.6),
         "Three DSP-anchored hyperparameters define the design",
         size=18, italic=False, color=RGBColor(0xCA, 0xDC, 0xFC), font=HDR_FONT)

lines = [
    ("n_heads", "≈  effective AR order"),
    ("context_length", "≈  N_FFT  ≈  fs / f_min"),
    ("head_dim", "≈  2 · log₂(context_length)"),
]
y = Inches(3.1)
for name, rule in lines:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y,
                              Inches(11.3), Inches(0.9))
    box.fill.solid(); box.fill.fore_color.rgb = DEEP; box.line.fill.background()
    add_text(s, Inches(1.3), y + Inches(0.2), Inches(3.5), Inches(0.5),
             name, size=20, bold=True, color=ACCENT, font=BODY_FONT)
    add_text(s, Inches(5.0), y + Inches(0.2), Inches(7), Inches(0.5),
             rule, size=20, color=LIGHT, font=BODY_FONT)
    y += Inches(1.05)

add_text(s, Inches(1.0), Inches(6.5), Inches(11.5), Inches(0.7),
         "emb_dim falls out.  The framing generalizes to many signal-processing tasks.",
         size=16, italic=False, color=RGBColor(0xCA, 0xDC, 0xFC), font=HDR_FONT)

out = r"C:\Users\Makarand Kulkarni\LLMs-from-scratch\.claude\worktrees\nice-banach\sinusoid_recovery\sinusoid_transformer_design.pptx"
prs.save(out)
print(f"Saved: {out}")
