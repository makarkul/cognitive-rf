"""Build sinusoid_transformer_design.pptx — 16-slide design walkthrough."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Ocean Gradient palette
DEEP = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
MID  = RGBColor(0x21, 0x29, 0x5C)
BG   = RGBColor(0xF5, 0xF7, 0xFA)
INK  = RGBColor(0x1A, 0x1A, 0x1A)
MUTE = RGBColor(0x55, 0x5F, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    return bg


def tb(slide, x, y, w, h, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.1))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = DEEP
    tb(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
       title, size=26, bold=True, color=WHITE)
    if subtitle:
        tb(slide, Inches(0.5), Inches(0.72), Inches(12), Inches(0.35),
           subtitle, size=13, color=RGBColor(0xCA, 0xDC, 0xFC))


def add_table(slide, x, y, w, h, data, header=True, col_widths=None):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT
            run.font.size = Pt(13)
            if header and r == 0:
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TEAL
            else:
                run.font.color.rgb = INK
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xE8, 0xEF, 0xF4)
    return tbl


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK)
add_bg(s, MID)
tb(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.4),
   "Sinusoid Recovery with a Tiny Transformer",
   size=40, bold=True, color=WHITE)
tb(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.6),
   "A DSP-to-Transformer Walkthrough",
   size=22, color=RGBColor(0xCA, 0xDC, 0xFC))
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.7),
                            Inches(1.5), Inches(0.06))
accent.line.fill.background()
accent.fill.solid()
accent.fill.fore_color.rgb = TEAL

# ---------- Slide 2: Motivation ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Motivation", "From classical DSP to attention-based prediction")
tb(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
   "Classical approach", size=18, bold=True, color=DEEP)
tb(s, Inches(0.8), Inches(2.0), Inches(12), Inches(2.0),
   "• Fit AR / MA models to a time series\n"
   "• Partial autocorrelation function (PACF) determines the model order\n"
   "• Well-understood, linear, minimum-variance under Gaussian assumptions",
   size=16)
tb(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
   "Thesis", size=18, bold=True, color=DEEP)
tb(s, Inches(0.8), Inches(4.7), Inches(12), Inches(2.0),
   "A transformer trained on next-sample prediction should implicitly\n"
   "learn the same AR structure — attention discovers which lags matter,\n"
   "the same way PACF does.",
   size=16)

# ---------- Slide 3: Key DSP fact ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Key DSP Fact", "A sinusoid is exactly AR(2)")
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.6))
box.line.fill.background(); box.fill.solid(); box.fill.fore_color.rgb = DEEP
tb(s, Inches(1.5), Inches(2.55), Inches(10.3), Inches(1.0),
   "x[n] = 2 cos(ω₀) · x[n-1]  −  x[n-2]",
   size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(4.3), Inches(12), Inches(2.5),
   "• Exact, not an approximation — any pure sinusoid obeys this recurrence\n"
   "• Only two effective taps are needed\n"
   "• Implication: the model should learn to attend primarily to lag 1 and lag 2",
   size=17)

# ---------- Slide 4: Problem setup ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Problem Setup", "Noisy sinusoid → predict the next clean sample")
data = [
    ["Parameter", "Range / Value"],
    ["Frequency f₀", "Uniform(1, 20) Hz"],
    ["Amplitude A", "Uniform(0.5, 2.0)"],
    ["Phase φ", "Uniform(0, 2π)"],
    ["Amplitude noise", "Multiplicative Gaussian"],
    ["Phase noise", "Additive Gaussian on φ"],
    ["Sample rate fs", "100 Hz"],
    ["Target", "Next clean sample from noisy context"],
]
add_table(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.6), data,
          col_widths=[Inches(4.5), Inches(7.2)])

# ---------- Slide 5: Architecture ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Architecture", "~29K parameters, 2 transformer blocks")
stages = [
    ("scalar\nsample", TEAL),
    ("Linear(1, 32)\ninput proj", DEEP),
    ("+ positional\nembedding", TEAL),
    ("2 × Transformer\nBlocks", DEEP),
    ("LayerNorm", TEAL),
    ("Linear(32, 1)", DEEP),
    ("predicted\nnext sample", TEAL),
]
x = Inches(0.3); y = Inches(3.0); w = Inches(1.72); h = Inches(1.4); gap = Inches(0.15)
for i, (label, col) in enumerate(stages):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.line.fill.background(); shp.fill.solid(); shp.fill.fore_color.rgb = col
    tb(s, x, y + Inches(0.35), w, Inches(0.8), label,
       size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + w, y + Inches(0.55), gap, Inches(0.3))
        arr.line.fill.background(); arr.fill.solid(); arr.fill.fore_color.rgb = MUTE
    x += w + gap
tb(s, Inches(0.6), Inches(5.2), Inches(12), Inches(1.5),
   "Total parameters: ~29K\n"
   "Input is a scalar (not a token); projection lifts it into a 32-D space\n"
   "where positional embeddings and attention can operate.",
   size=15)

# ---------- Slide 6: TL;DR design table ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "The TL;DR Design Table", "Three DSP-anchored hyperparameters")
data = [
    ["Hyperparameter", "DSP analogue", "Value for this task"],
    ["n_heads", "AR order", "2  (sinusoid is AR(2))"],
    ["context_length", "≈ N_FFT = fs / f_min", "100 / 1 ≈ 128"],
    ["head_dim", "codeword length ≈ 2·log₂(context)", "2·log₂(128) = 16"],
]
add_table(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(3.0), data,
          col_widths=[Inches(3.0), Inches(5.0), Inches(4.1)])
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.6))
box.line.fill.background(); box.fill.solid(); box.fill.fore_color.rgb = DEEP
tb(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.2),
   "emb_dim = n_heads × head_dim = 2 × 16 = 32   (forced by construction)",
   size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ---------- Slide 7: Tokenization ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Tokenization Is Different", "Continuous reals, not discrete tokens")
tb(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.5),
   "Text LLM", size=18, bold=True, color=DEEP)
tb(s, Inches(0.8), Inches(2.1), Inches(5.6), Inches(2.5),
   "• Tokenizer splits text into IDs\n"
   "• vocab_size entries\n"
   "• nn.Embedding(vocab_size, 32)\n"
   "• Discrete input",
   size=15)
tb(s, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.5),
   "This model", size=18, bold=True, color=TEAL)
tb(s, Inches(7.2), Inches(2.1), Inches(5.6), Inches(2.5),
   "• No tokenizer\n"
   "• No vocab_size\n"
   "• Linear(1, 32) lifts each scalar\n"
   "• Continuous real input",
   size=15)
tb(s, Inches(0.6), Inches(5.2), Inches(12.2), Inches(1.8),
   "Each sample x[n] ∈ ℝ becomes a 32-D embedding through a linear projection.\n"
   "This is the only structural change vs. a GPT-style block.",
   size=15, color=MUTE)

# ---------- Slide 8: Is 32-D wasted ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Is the 32-D Embedding Wasted?", "Short answer: no — but not for the scalar")
tb(s, Inches(0.6), Inches(1.6), Inches(12), Inches(2.2),
   "Linear(1, 32) output is rank-1 in the sample value v:\n"
   "    emb(v) = v · w   where w ∈ ℝ³²\n"
   "All 32 dims are perfectly correlated — they carry exactly one number.",
   size=16)
tb(s, Inches(0.6), Inches(4.0), Inches(12), Inches(2.8),
   "So why 32 dims?\n"
   "• Positional embeddings live in the same 32-D space\n"
   "• Multi-head attention needs room to split into n_heads × head_dim\n"
   "• Residual stream accumulates features across layers\n"
   "The 32-D space exists for positional + attention machinery, not the scalar itself.",
   size=15)

# ---------- Slide 9: context_length <-> N_FFT ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "context_length ↔ N_FFT", "Cover the slowest period in the signal")
tb(s, Inches(0.6), Inches(1.5), Inches(12), Inches(1.8),
   "Rule: context_length ≈ fs / f_min\n"
   "Here: 100 Hz / 1 Hz = 100 samples → rounded to 128 (power of 2).",
   size=17)
data = [
    ["LTE bandwidth", "N_FFT", "Comment"],
    ["1.4 MHz", "128", "smallest"],
    ["3 MHz", "256", ""],
    ["5 MHz", "512", ""],
    ["10 MHz", "1024", ""],
    ["20 MHz", "2048", "largest"],
]
add_table(s, Inches(2.5), Inches(3.4), Inches(8.3), Inches(3.6), data,
          col_widths=[Inches(2.6), Inches(2.0), Inches(3.7)])
tb(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
   "Invariant: N_FFT / fs stays roughly constant across configurations.",
   size=14, color=MUTE, align=PP_ALIGN.CENTER)

# ---------- Slide 10: n_heads ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "n_heads = Number of Lags That Matter", "One head per relevant AR tap")
tb(s, Inches(0.6), Inches(1.7), Inches(12), Inches(5),
   "For an AR(p) process, you want roughly p heads —\n"
   "each head can specialize on a distinct lag pattern.\n\n"
   "Sinusoid → AR(2) → n_heads = 2\n"
   "    Head 0 learns to look at lag 1\n"
   "    Head 1 learns to look at lag 2\n\n"
   "General heuristic:   n_heads ≈ effective AR order of the signal.",
   size=17)

# ---------- Slide 10b: n_layers rationale ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Why n_layers = 2", "Depth = number of sequential reasoning steps")
tb(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(5.5),
   "ω₀ is unknown and varies per example (1–20 Hz).\n"
   "The model cannot apply AR(2) until it first estimates ω₀.\n"
   "That is a two-step computation — the minimum depth is 2.\n\n"
   "    Layer 1:  estimate ω₀ from the context  (frequency estimation)\n"
   "    Layer 2:  apply the AR(2) recurrence using ω₀  (filtering)\n\n"
   "General rule:  n_layers ≈ number of sequential sub-tasks the problem requires.\n\n"
   "More layers add no new structure for pure sinusoids —\n"
   "empirically, deeper models did not improve sinusoid recovery.",
   size=16)

# ---------- Slide 10c: heads per layer ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Two Heads Per Layer — Different Jobs", "Same count, different roles")
data = [
    ["Layer", "Role", "Why 2 heads"],
    ["Layer 1", "Estimate ω₀",
     "Need two projections / lags to disambiguate frequency\n"
     "(e.g. sine + cosine bases, or lag-1 + lag-2 autocorrelation)"],
    ["Layer 2", "Apply AR(2) filter",
     "One head per filter tap: lag 1 with weight 2cos(ω₀),\n"
     "lag 2 with weight −1"],
]
add_table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.2), data,
          col_widths=[Inches(1.6), Inches(3.2), Inches(7.5)])
tb(s, Inches(0.6), Inches(5.1), Inches(12.2), Inches(2.2),
   "Refined rule:\n"
   "    n_heads  ≥  max( AR order,  intrinsic rank of the signal in embedding space )\n\n"
   "For a pure sinusoid both are 2, so the same head count serves both layers —\n"
   "but the reason is different in each layer.",
   size=15)

# ---------- Slide 11: head_dim as channel coding ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "head_dim as Channel Coding", "Dimension of the Q·K dot product")
tb(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(5.5),
   "head_dim governs how many positions attention can discriminate robustly.\n\n"
   "Classical bounds point to the same sizing:\n"
   "    • Shannon capacity of a noisy channel\n"
   "    • Spherical codes on the unit sphere Sᵈ⁻¹\n"
   "    • BCH / Reed–Muller codeword length\n\n"
   "All suggest:   head_dim  ~  2 · log₂(context_length)",
   size=16)

# ---------- Slide 12: Channel coding analogy ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Channel Coding Analogy", "Attention = matched filtering")
tb(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(5.5),
   "• Q and K behave like spreading codes / matched-filter templates\n"
   "• The  1 / √head_dim  scaling in attention is exactly\n"
   "  the matched-filter processing gain from communications theory\n"
   "• Longer head_dim  →  lower cross-correlation between distinct patterns\n"
   "  →  less interference between attention queries\n\n"
   "Short head_dim behaves like a short spreading code: noisy, ambiguous matches.",
   size=16)

# ---------- Slide 13: Shannon sizing check ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Shannon Sizing Check", "Does head_dim = 16 have enough room?")
data = [
    ["Quantity", "Value"],
    ["Positions to discriminate", "context_length = 128"],
    ["Bits of position info needed", "log₂(128) = 7"],
    ["Per-dim SNR (post-softmax)", "~ 1"],
    ["Min dims for 7 bits @ SNR 1", "d ≥ 14"],
    ["Our head_dim", "16  (≈ 1× safety factor)"],
    ["Typical LLM safety factor", "3× – 5×  (harder tasks)"],
]
add_table(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(4.5), data,
          col_widths=[Inches(5.5), Inches(5.8)])
tb(s, Inches(0.6), Inches(6.3), Inches(12.2), Inches(0.8),
   "head_dim = 16 lands right at the Shannon floor — sufficient for this simple task.",
   size=14, color=MUTE, align=PP_ALIGN.CENTER)

# ---------- Slide 14: Training ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Training", "MSE, AdamW, 30 epochs on CPU")
tb(s, Inches(0.6), Inches(1.6), Inches(12), Inches(5.5),
   "Loss:          MSE between predicted and clean next sample\n"
   "Optimizer:     AdamW\n"
   "Epochs:        30\n"
   "Hardware:      CPU (tiny model)\n\n"
   "Result:\n"
   "    • Training MSE:    1.4  →  0.03\n"
   "    • Output SNR:      ~ 14.8 dB\n"
   "    • Matches the naive denoising baseline\n\n"
   "Interpretation: the model predicts the next CLEAN sample\n"
   "from noisy history — it has learned to denoise while forecasting.",
   size=16)

# ---------- Slide 15: Analysis diagnostics ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "Analysis Diagnostics", "The model behaves like AR(2)")
tb(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(5.5),
   "Evidence across four views:\n\n"
   "1.  Attention heatmaps — peaks concentrated at lag 1 and lag 2\n\n"
   "2.  Attention-vs-lag profile — matches AR(2) theoretical taps\n\n"
   "3.  FFT of the output — noise floor suppressed, peak at f₀ preserved\n\n"
   "4.  Autoregressive generation — produces clean sinusoids when seeded\n"
   "    with a noisy prefix",
   size=16)

# ---------- Slide 16: Takeaway ----------
s = prs.slides.add_slide(BLANK); add_bg(s)
add_bg(s, MID)
tb(s, Inches(0.8), Inches(0.6), Inches(12), Inches(0.8),
   "Key Takeaway", size=30, bold=True, color=WHITE)
tb(s, Inches(0.8), Inches(1.6), Inches(12), Inches(5.5),
   "Three DSP-anchored hyperparameters drive the entire design:\n\n"
   "    n_heads         =  AR order of the signal\n"
   "    context_length  =  N_FFT  =  fs / f_min\n"
   "    head_dim        =  2 · log₂(context_length)\n\n"
   "emb_dim falls out automatically as  n_heads × head_dim.\n\n"
   "This framing generalizes to many signal-processing tasks —\n"
   "the transformer is a learned, data-driven AR model whose\n"
   "sizing is governed by classical DSP and channel-coding bounds.",
   size=17, color=WHITE)

out = r"C:\Users\Makarand Kulkarni\LLMs-from-scratch\.claude\worktrees\nice-banach\sinusoid_recovery\sinusoid_transformer_design.pptx"
prs.save(out)
print(f"Saved: {out}")
